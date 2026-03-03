#!/usr/bin/env python3
"""
pptx_builder.py — 템플릿 기반 PPTX 생성기

Usage:
  python pptx_builder.py --slides slides.json [options]

Options:
  --slides SLIDES_JSON      슬라이드 데이터 JSON 파일 (필수)
  --template TEMPLATE       빌트인 테마: business|minimal|technical|pitch (기본: business)
  --custom-theme THEME_JSON 커스텀 테마 JSON 파일 (빌트인 테마 위에 덮어씀)
  --output OUTPUT           출력 파일 경로 (기본: output.pptx)

Slides JSON 구조:
  [
    {"type": "title", "title": "제목", "subtitle": "부제목"},
    {"type": "content", "title": "슬라이드 제목", "bullets": ["항목1", "항목2"]},
    {"type": "two_col", "title": "비교", "left": ["A", "B"], "right": ["C", "D"]},
    {"type": "table", "title": "데이터", "headers": ["A","B"], "rows": [["1","2"]]},
    {"type": "image", "title": "이미지", "path": "chart.png"},
    {"type": "blank", "title": "빈 슬라이드"}
  ]

Custom theme JSON 구조 (모든 필드 선택적):
  {
    "bg_dark":        "1a1a2e",   # 배경(진) 헥스
    "bg_light":       "162132",   # 배경(연) 헥스
    "accent":         "0f3f5c",   # 강조 색 헥스
    "text_primary":   "ffffff",   # 기본 텍스트 헥스
    "text_secondary": "aabbcc",   # 보조 텍스트 헥스
    "font_title":     "Arial",    # 제목 폰트명
    "font_body":      "Calibri"   # 본문 폰트명
  }
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 빌트인 테마 ─────────────────────────────────────────────
BUILTIN_THEMES: dict[str, dict] = {
    "business": {
        "bg_dark":        "1a1a2e",
        "bg_light":       "16213e",
        "accent":         "0f3f5c",
        "text_primary":   "ffffff",
        "text_secondary": "aabbcc",
        "font_title":     "Calibri",
        "font_body":      "Calibri",
    },
    "minimal": {
        "bg_dark":        "f5f5f5",
        "bg_light":       "ffffff",
        "accent":         "333333",
        "text_primary":   "111111",
        "text_secondary": "666666",
        "font_title":     "Arial",
        "font_body":      "Arial",
    },
    "technical": {
        "bg_dark":        "1e1e1e",
        "bg_light":       "252526",
        "accent":         "007acc",
        "text_primary":   "d4d4d4",
        "text_secondary": "858585",
        "font_title":     "Consolas",
        "font_body":      "Consolas",
    },
    "pitch": {
        "bg_dark":        "0d0d0d",
        "bg_light":       "1a1a1a",
        "accent":         "c9a84c",
        "text_primary":   "ffffff",
        "text_secondary": "ccbbaa",
        "font_title":     "Georgia",
        "font_body":      "Calibri",
    },
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_theme(template: str, custom_path: str | None) -> dict:
    theme = dict(BUILTIN_THEMES.get(template, BUILTIN_THEMES["business"]))
    if custom_path:
        overrides = json.loads(Path(custom_path).read_text(encoding="utf-8"))
        theme.update(overrides)
    return theme


# ── 슬라이드 빌더 ────────────────────────────────────────────
def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height,
                 text, font_name, font_size, bold, color, align=PP_ALIGN.LEFT) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_title_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_dark"]))

    w, h = prs.slide_width, prs.slide_height
    # 가운데 강조 바
    bar = slide.shapes.add_shape(1, 0, int(h * 0.4), w, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    title_text = data.get("title", "")
    subtitle_text = data.get("subtitle", "")

    _add_textbox(slide, Inches(1), Inches(2), Inches(11.33), Inches(1.5),
                 title_text, theme["font_title"], 44, True,
                 hex_to_rgb(theme["text_primary"]), PP_ALIGN.CENTER)

    if subtitle_text:
        _add_textbox(slide, Inches(1), Inches(3.8), Inches(11.33), Inches(1),
                     subtitle_text, theme["font_body"], 22, False,
                     hex_to_rgb(theme["text_secondary"]), PP_ALIGN.CENTER)


def build_content_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_light"]))

    w = prs.slide_width
    # 제목 바
    bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 data.get("title", ""), theme["font_title"], 28, True,
                 hex_to_rgb(theme["text_primary"]))

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(data.get("bullets", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.name = theme["font_body"]
        run.font.size = Pt(20)
        run.font.color.rgb = hex_to_rgb(theme["text_primary"])


def build_two_col_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_light"]))

    w = prs.slide_width
    bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 data.get("title", ""), theme["font_title"], 28, True,
                 hex_to_rgb(theme["text_primary"]))

    # 구분선
    div = slide.shapes.add_shape(1, int(w / 2) - Inches(0.03), Inches(1.3), Inches(0.06), Inches(5.8))
    div.fill.solid()
    div.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    div.line.fill.background()

    for side, x_start in [("left", 0.5), ("right", 6.9)]:
        items = data.get(side, [])
        body = slide.shapes.add_textbox(Inches(x_start), Inches(1.5), Inches(5.5), Inches(5.5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = theme["font_body"]
            run.font.size = Pt(18)
            run.font.color.rgb = hex_to_rgb(theme["text_primary"])


def build_table_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_light"]))

    w = prs.slide_width
    bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 data.get("title", ""), theme["font_title"], 28, True,
                 hex_to_rgb(theme["text_primary"]))

    headers = data.get("headers", [])
    rows_data = data.get("rows", [])
    if not headers:
        return

    cols = len(headers)
    rows = len(rows_data) + 1
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5 * rows))
    tbl = tbl_shape.table

    accent = hex_to_rgb(theme["accent"])
    primary = hex_to_rgb(theme["text_primary"])
    bg_light = hex_to_rgb(theme["bg_dark"])

    for col_i, header in enumerate(headers):
        cell = tbl.cell(0, col_i)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.bold = True
        run.font.name = theme["font_title"]
        run.font.size = Pt(14)
        run.font.color.rgb = primary
        cell.fill.solid()
        cell.fill.fore_color.rgb = accent

    for row_i, row in enumerate(rows_data, start=1):
        for col_i, val in enumerate(row[:cols]):
            cell = tbl.cell(row_i, col_i)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = theme["font_body"]
            run.font.size = Pt(12)
            run.font.color.rgb = primary
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_light


def build_image_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_light"]))

    w = prs.slide_width
    bar = slide.shapes.add_shape(1, 0, 0, w, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(theme["accent"])
    bar.line.fill.background()

    _add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
                 data.get("title", ""), theme["font_title"], 28, True,
                 hex_to_rgb(theme["text_primary"]))

    img_path = data.get("path", "")
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.4), Inches(10), Inches(5.5))
    else:
        _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                     f"[이미지 없음: {img_path}]", theme["font_body"], 16, False,
                     hex_to_rgb(theme["text_secondary"]), PP_ALIGN.CENTER)


def build_blank_slide(prs: Presentation, data: dict, t: dict, theme: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, hex_to_rgb(theme["bg_light"]))
    if data.get("title"):
        _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
                     data["title"], theme["font_title"], 28, True,
                     hex_to_rgb(theme["text_primary"]))


SLIDE_BUILDERS = {
    "title":   build_title_slide,
    "content": build_content_slide,
    "two_col": build_two_col_slide,
    "table":   build_table_slide,
    "image":   build_image_slide,
    "blank":   build_blank_slide,
}


# ── 메인 ─────────────────────────────────────────────────────
def main() -> None:
    parser = argparse.ArgumentParser(description="템플릿 기반 PPTX 생성기")
    parser.add_argument("--slides",        required=True, help="슬라이드 JSON 파일")
    parser.add_argument("--template",      default="business",
                        choices=list(BUILTIN_THEMES), help="빌트인 테마")
    parser.add_argument("--custom-theme",  default=None, help="커스텀 테마 JSON 파일")
    parser.add_argument("--output",        default="output.pptx", help="출력 파일")
    args = parser.parse_args()

    slides_data: list[dict] = json.loads(Path(args.slides).read_text(encoding="utf-8"))
    theme = load_theme(args.template, args.custom_theme)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type)
        if builder is None:
            print(f"  [skip] 알 수 없는 슬라이드 타입: {slide_type} (슬라이드 #{i+1})", file=sys.stderr)
            continue
        builder(prs, slide_data, {}, theme)
        print(f"  [{i+1}] {slide_type}: {slide_data.get('title', '(제목 없음)')}")

    prs.save(args.output)
    print(f"saved: {args.output}")


if __name__ == "__main__":
    main()
